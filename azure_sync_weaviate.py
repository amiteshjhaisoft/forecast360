# Author: Amitesh Jha | iSoft

from __future__ import annotations
import io, os, json, hashlib, re
from typing import Any, Dict, List, Optional, Iterable, Tuple

import numpy as np
from azure.identity import DefaultAzureCredential
from azure.storage.blob import BlobServiceClient, ContentSettings

import weaviate
import weaviate.classes as wvc

# -------- Optional extractors (all graceful) --------
# PDF
try:
    import pypdf  # uncomment in requirements if you want real PDF text
except Exception:
    pypdf = None
# DOCX
try:
    import docx2txt  # uncomment in requirements to enable DOCX text
except Exception:
    docx2txt = None
# PPTX
try:
    from pptx import Presentation  # python-pptx
except Exception:
    Presentation = None
# HTML tidy (optional)
try:
    from bs4 import BeautifulSoup  # beautifulsoup4 + (optionally) lxml
except Exception:
    BeautifulSoup = None
# XLSX
try:
    import openpyxl  # already in your requirements
except Exception:
    openpyxl = None
# Images (EXIF)
try:
    from PIL import Image, ExifTags  # already in your requirements
    _EXIF_TAGS = {v: k for k, v in ExifTags.TAGS.items()}
except Exception:
    Image, ExifTags, _EXIF_TAGS = None, None, {}

# Optional local embeddings (used when vectorizer='none')
try:
    from sentence_transformers import SentenceTransformer
except Exception:
    SentenceTransformer = None  # graceful


# ----------------------- Secrets helpers -----------------------
def _sget(st_secrets, section: str, key: str, default: Any = None) -> Any:
    try:
        return st_secrets[section].get(key, default)  # type: ignore
    except Exception:
        return default


# ----------------------- Extension groups -----------------------
_TEXT_EXT  = {".txt", ".md", ".json", ".csv", ".tsv", ".html", ".htm"}
_PDF_EXT   = {".pdf"}
_DOC_EXT   = {".docx"}  # (legacy .doc not supported without extra deps)
_PPTX_EXT  = {".pptx"}
_XLSX_EXT  = {".xlsx"}  # (.xls handled only if xlrd==1.2.0 + separate path)
_IMG_EXT   = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tiff", ".tif", ".gif"}
_AUDIO_EXT = {".mp3", ".wav", ".m4a", ".ogg", ".flac", ".aac"}
_VIDEO_EXT = {".mp4", ".mkv", ".mov", ".avi", ".wmv", ".webm", ".m4v"}

# If a blob is larger than this (bytes), avoid downloading entire content for text extraction.
# You can raise this if you want bigger files extracted.
MAX_DOWNLOAD_BYTES = 15 * 1024 * 1024  # 15 MB


def _ext(path: str) -> str:
    m = re.search(r"(\.[A-Za-z0-9]+)$", path)
    return m.group(1).lower() if m else ""


def _sha1(data: bytes) -> str:
    return hashlib.sha1(data).hexdigest()


# ----------------------- Basic parsers -----------------------
def _clean_html_to_text(html: str) -> str:
    if BeautifulSoup:
        soup = BeautifulSoup(html, "lxml") if "lxml" else BeautifulSoup(html, "html.parser")
        for t in soup(["script", "style"]):
            t.decompose()
        raw = soup.get_text(" ", strip=True)
    else:
        raw = re.sub(r"<(script|style)[\s\S]*?</\1>", " ", html, flags=re.I)
        raw = re.sub(r"<[^>]+>", " ", raw)
    raw = re.sub(r"\s+", " ", raw).strip()
    return raw


def _bytes_to_text_guess(name: str, data: bytes, content_type: Optional[str] = None) -> Optional[str]:
    """
    Try to decode as UTF-8; fallback to ISO-8859-1; else None.
    """
    try:
        return data.decode("utf-8", errors="strict")
    except Exception:
        try:
            return data.decode("utf-8", errors="ignore")
        except Exception:
            try:
                return data.decode("latin-1", errors="ignore")
            except Exception:
                return None


def _xlsx_to_text(data: bytes) -> Optional[str]:
    if not openpyxl:
        return None
    try:
        f = io.BytesIO(data)
        wb = openpyxl.load_workbook(f, read_only=True, data_only=True)
        parts: List[str] = []
        for ws in wb.worksheets[:3]:  # sample first 3 sheets
            parts.append(f"[Sheet] {ws.title}")
            rows = ws.iter_rows(min_row=1, max_row=50, values_only=True)  # sample top-50 rows
            for r in rows:
                row_txt = " | ".join("" if v is None else str(v) for v in r)
                if row_txt.strip():
                    parts.append(row_txt)
        return "\n".join(parts).strip()[:150_000]  # limit size
    except Exception:
        return None


def _pdf_to_text(data: bytes) -> Optional[str]:
    if not pypdf:
        return None
    try:
        reader = pypdf.PdfReader(io.BytesIO(data))
        pages = []
        for i, pg in enumerate(reader.pages):
            if i >= 50:  # cap pages for huge PDFs
                break
            try:
                pages.append(pg.extract_text() or "")
            except Exception:
                continue
        txt = "\n".join(pages).strip()
        return txt or None
    except Exception:
        return None


def _docx_to_text(data: bytes) -> Optional[str]:
    if not docx2txt:
        return None
    try:
        # docx2txt requires a path; write to temp buffer
        bio = io.BytesIO(data)
        tmp = io.BytesIO()
        # workaround: docx2txt doesn't handle BytesIO directly; write to temp file on disk
        import tempfile
        with tempfile.NamedTemporaryFile(suffix=".docx", delete=True) as tf:
            tf.write(data); tf.flush()
            text = docx2txt.process(tf.name) or ""
        return text.strip() or None
    except Exception:
        return None


def _pptx_to_text(data: bytes) -> Optional[str]:
    if not Presentation:
        return None
    try:
        prs = Presentation(io.BytesIO(data))
        chunks: List[str] = []
        for i, slide in enumerate(prs.slides):
            if i >= 200:
                break
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    t = (shape.text or "").strip()
                    if t:
                        texts.append(t)
            if texts:
                chunks.append(f"[Slide {i+1}] " + " | ".join(texts))
        return "\n".join(chunks).strip() or None
    except Exception:
        return None


def _image_to_surrogate(name: str, size: int, data: Optional[bytes]) -> str:
    if Image and data:
        try:
            im = Image.open(io.BytesIO(data))
            w, h = im.size
            fmt = im.format
            meta = f"Image {fmt} {w}x{h}"
            try:
                exif = im._getexif() or {}
                cam = exif.get(_EXIF_TAGS.get("Model", -1), None)
                dt  = exif.get(_EXIF_TAGS.get("DateTime", -1), None)
                if cam or dt:
                    meta += f" | EXIF: " + " ".join([x for x in [str(cam), str(dt)] if x and x != 'None'])
            except Exception:
                pass
            return f"[IMAGE] {name} ({size} bytes) — {meta}"
        except Exception:
            pass
    return f"[IMAGE] {name} ({size} bytes)"


def _binary_surrogate(kind: str, name: str, size: int, mime: Optional[str]) -> str:
    m = f"{mime}" if mime else "application/octet-stream"
    return f"[{kind.upper()}] {name} ({size} bytes) — {m}"


# ----------------------- Text extraction (multi-type) -----------------------
def _to_text(name: str, data: bytes, content_type: Optional[str]) -> Optional[str]:
    ext = _ext(name)

    # Plain text-ish
    if ext in {".txt", ".md", ".json", ".csv", ".tsv"}:
        return _bytes_to_text_guess(name, data, content_type)

    # HTML
    if ext in {".html", ".htm"}:
        html = _bytes_to_text_guess(name, data, content_type) or ""
        return _clean_html_to_text(html)

    # PDFs
    if ext in _PDF_EXT:
        return _pdf_to_text(data)

    # DOCX
    if ext in _DOC_EXT:
        return _docx_to_text(data)

    # PPTX
    if ext in _PPTX_EXT:
        return _pptx_to_text(data)

    # XLSX
    if ext in _XLSX_EXT:
        return _xlsx_to_text(data)

    # Otherwise we don't have a parser
    return None


# ----------------------- Chunking -----------------------
def _chunk_text(text: str, max_chars: int = 1400, overlap: int = 200) -> List[str]:
    text = re.sub(r"\s+", " ", text).strip()
    if not text:
        return []
    chunks: List[str] = []
    i = 0
    n = len(text)
    while i < n:
        end = min(n, i + max_chars)
        # try to cut at sentence boundary
        cut = re.search(r"(?<=\.|\?|!)\s", text[i:end])
        if cut:
            end = i + cut.end()
        chunk = text[i:end].strip()
        if chunk:
            chunks.append(chunk)
        i = end - overlap if end - overlap > i else end
    return chunks


# ----------------------- Embedding (optional) -----------------------
def _maybe_embedder(model_name: Optional[str]) -> Optional["SentenceTransformer"]:
    if not model_name or SentenceTransformer is None:
        return None
    try:
        return SentenceTransformer(model_name)
    except Exception:
        return None


# ----------------------- Azure Blob helpers -----------------------
def _azure_blob_client(st_secrets) -> BlobServiceClient:
    conn_str = _sget(st_secrets, "azure", "connection_string") or _sget(st_secrets, "azure_blob", "connection_string")
    if conn_str:
        return BlobServiceClient.from_connection_string(conn_str)

    account_url = _sget(st_secrets, "azure", "account_url") or _sget(st_secrets, "azure_blob", "account_url")
    if not account_url:
        raise RuntimeError("Azure Blob config missing. Provide either connection_string or account_url in secrets.")
    cred = DefaultAzureCredential(exclude_interactive_browser_credential=False)
    return BlobServiceClient(account_url=account_url, credential=cred)


# ----------------------- Weaviate helpers -----------------------
def _weaviate_client(st_secrets) -> weaviate.WeaviateClient:
    url = _sget(st_secrets, "weaviate", "url", "")
    if not url:
        raise RuntimeError("Set [weaviate].url in secrets.")
    api_key = _sget(st_secrets, "weaviate", "api_key", "")
    auth = wvc.init.Auth.api_key(api_key) if api_key else None

    client = weaviate.connect_to_weaviate_cloud(
        cluster_url=url,
        auth_credentials=auth,
        additional_config=wvc.init.AdditionalConfig(
            timeout=wvc.init.Timeout(init=30, query=60, insert=120)
        ),
    )
    return client


def _collection_info(coll) -> Tuple[bool, Optional[int]]:
    """Return (needs_manual_vectors, dims_if_known) from collection config."""
    try:
        cfg = coll.config.get()
        vcfg = getattr(cfg, "vectorizer", None)
        needs_vec = (str(vcfg).lower() == "none")
        return needs_vec, None
    except Exception:
        return False, None


# ----------------------- Core sync -----------------------
def sync_from_azure(
    st_secrets,
    collection_name: str,
    container_key: str = "container",     # secrets['azure' or 'azure_blob'][container_key]
    prefix_key: str = "prefix",           # optional subfolder in the container
    embed_model_key: str = ("rag", "embed_model"),  # ('section','key') for local embedding model
    delete_before_upsert: bool = True,
    max_docs: Optional[int] = None
) -> Dict[str, Any]:
    """
    Reads files from an Azure Blob *folder* and upserts chunks/surrogates into a Weaviate v4 collection.
    Properties stored: text, source
    - If vectorizer=none, we compute vectors locally over the text/surrogate and send them.
    - For binaries (images/audio/video/other) we insert a textual surrogate so they're still searchable.
    """
    container = _sget(st_secrets, "azure", container_key) or _sget(st_secrets, "azure_blob", container_key)
    if not container:
        raise RuntimeError("Set [azure].[container] (or [azure_blob]) in secrets.")
    prefix = _sget(st_secrets, "azure", prefix_key) or _sget(st_secrets, "azure_blob", prefix_key) or ""

    # Connect
    bs = _azure_blob_client(st_secrets)
    client = _weaviate_client(st_secrets)
    coll = client.collections.use(collection_name)

    needs_vec, _ = _collection_info(coll)
    embed_model_name = None
    if isinstance(embed_model_key, (tuple, list)) and len(embed_model_key) == 2:
        embed_model_name = _sget(st_secrets, embed_model_key[0], embed_model_key[1])
    embedder = _maybe_embedder(embed_model_name) if needs_vec else None

    container_client = bs.get_container_client(container)
    blob_iter = container_client.list_blobs(name_starts_with=prefix or None)

    inserted, skipped, cleared_sources = 0, 0, 0
    processed_docs = 0

    batch_props: List[Dict[str, Any]] = []
    batch_vecs: Optional[List[List[float]]] = [] if needs_vec else None
    batch_size = 64

    def _flush():
        nonlocal inserted, batch_props, batch_vecs
        if not batch_props:
            return
        if needs_vec:
            coll.data.insert_many(objects=batch_props, vectors=batch_vecs)  # vectors aligned to objects
        else:
            coll.data.insert_many(objects=batch_props)
        inserted += len(batch_props)
        batch_props = []
        if batch_vecs is not None:
            batch_vecs = []

    for blob in blob_iter:
        if max_docs and processed_docs >= max_docs:
            break

        name = blob.name
        ext = _ext(name)
        size = int(getattr(blob, "size", 0) or 0)
        content_type = None
        try:
            content_type = (getattr(blob, "content_settings", None) or ContentSettings()).content_type
        except Exception:
            content_type = None

        src = f"azure://{container}/{name}"

        # Clear old records for this source (idempotent refresh)
        if delete_before_upsert:
            try:
                coll.data.delete_many(where=wvc.query.Filter.by_property("source").equal(src))
                cleared_sources += 1
            except Exception:
                pass

        # Decide whether to download body or create surrogate
        need_download = False
        is_text_like = ext in (_TEXT_EXT | _PDF_EXT | _DOC_EXT | _PPTX_EXT | _XLSX_EXT)
        is_image = ext in _IMG_EXT
        is_audio = ext in _AUDIO_EXT
        is_video = ext in _VIDEO_EXT

        # Download if text-like and not too big
        if is_text_like and (size == 0 or size <= MAX_DOWNLOAD_BYTES):
            need_download = True

        data: Optional[bytes] = None
        if need_download:
            try:
                bclient = container_client.get_blob_client(name)
                data = bclient.download_blob(max_concurrency=1).readall()
            except Exception:
                data = None

        # Extract or build surrogate
        extracted_text: Optional[str] = None
        if data is not None and is_text_like:
            extracted_text = _to_text(name, data, content_type)
        elif is_image:
            # For small images, try to open to get W×H
            img_bytes: Optional[bytes] = None
            if size and size <= 8 * 1024 * 1024:  # 8MB cap for image attributes
                try:
                    bclient = container_client.get_blob_client(name)
                    img_bytes = bclient.download_blob(max_concurrency=1).readall()
                except Exception:
                    img_bytes = None
            extracted_text = _image_to_surrogate(name, size, img_bytes)
        elif is_audio:
            extracted_text = _binary_surrogate("audio", name, size, content_type)
        elif is_video:
            extracted_text = _binary_surrogate("video", name, size, content_type)
        else:
            # other binaries
            extracted_text = _binary_surrogate("binary", name, size, content_type)

        # If nothing at all (e.g., huge text files we skipped), build minimal surrogate
        if not extracted_text:
            extracted_text = f"[FILE] {name} ({size} bytes)" + (f" — {content_type}" if content_type else "")

        # Chunk or single object
        chunks = _chunk_text(extracted_text, max_chars=1400, overlap=200) if len(extracted_text) > 1600 else [extracted_text]

        for ch in chunks:
            obj = {"text": ch, "source": src}
            batch_props.append(obj)
            if needs_vec:
                if embedder is None:
                    client.close()
                    raise RuntimeError(
                        "Collection is configured with vectorizer='none' but no local embedding model is available. "
                        "Set [rag].embed_model in secrets and install sentence-transformers."
                    )
                vec = embedder.encode([ch], normalize_embeddings=True)[0].astype("float32").tolist()
                batch_vecs.append(vec)  # type: ignore
            if len(batch_props) >= batch_size:
                _flush()

        processed_docs += 1

    _flush()
    client.close()
    return {
        "inserted_chunks": inserted,
        "skipped_files": skipped,  # only increments when we truly skip (now rare)
        "cleared_sources": cleared_sources,
        "processed_files": processed_docs,
        "vectorizer_none": needs_vec,
        "used_embed_model": embed_model_name if needs_vec else None,
    }
